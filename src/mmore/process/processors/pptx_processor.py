import io
import logging
from typing import cast

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.autoshape import Shape
from pptx.shapes.picture import Picture

from ...type import DocumentMetadata, FileDescriptor, MultimodalSample
from ..utils import clean_image, clean_text
from .base import Processor, ProcessorConfig

logger = logging.getLogger(__name__)


class PPTXProcessor(Processor):
    """
    A processor for handling PPTX files. Extracts text, images, and notes from PowerPoint presentations.

    Attributes:
        files (List[FileDescriptor]): List of files to be processed.
        config (ProcessorConfig): Configuration for the processor.
    """

    def __init__(self, config=None):
        """
        Args:
            files (List[FileDescriptor]): List of files to process.
            config (ProcessorConfig, optional): Configuration for the processor. Defaults to None.
        """
        super().__init__(config=config or ProcessorConfig())

    @classmethod
    def accepts(cls, file: FileDescriptor) -> bool:
        """
        Args:
            file (FileDescriptor): The file descriptor to check.

        Returns:
            bool: True if the file is a PPTX file, False otherwise.
        """
        return file.file_extension.lower() in [".pptx"]

    def process(self, file_path: str) -> MultimodalSample:
        """
        Process a single PPTX file. Extracts text, images, and notes from each slide.

        Args:
            file_path (str): Path to the PPTX file.

        Returns:
            dict: A dictionary containing processed text and images.

        The method processes each slide, extracting text and images from shapes,
        and extracts notes if present. The elements are sorted by their vertical position.
        """

        logger.info(f"Processing PowerPoint file: {file_path}")
        try:
            prs = Presentation(file_path)
        except Exception as e:
            logger.error(f"Failed to open PowerPoint file {file_path}: {e}")
            return self.create_sample([], [], DocumentMetadata(file_path=file_path))

        all_text: list[str] = []
        embedded_images: list[Image.Image] = []

        try:
            for slide in prs.slides:
                # 1) Extract text and images from slide
                # Sort shapes by their vertical position
                shape_list = sorted(
                    (shape for shape in slide.shapes if hasattr(shape, "top")),
                    key=lambda s: s.top,
                )

                for shape in shape_list:
                    # Extract text from shape
                    if shape.has_text_frame:
                        cleaned_text = clean_text(cast(Shape, shape).text)
                        if cleaned_text.strip():
                            all_text.append(cleaned_text)

                    # Extract images from shape
                    if self.config.custom_config.get("extract_images", True):
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                pil_image = Image.open(
                                    io.BytesIO(cast(Picture, shape).image.blob)
                                ).convert("RGBA")
                                if clean_image(pil_image):
                                    embedded_images.append(pil_image)
                                    all_text.append(self.config.attachment_tag)

                            except Exception as e:
                                logger.error(f"Failed to extract image from slide: {e}")
                    else:
                        embedded_images = []

                # 2) Extract text from slide notes if present
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame
                    for paragraph in notes.paragraphs:
                        if paragraph.text:
                            cleaned = clean_text(paragraph.text)
                            if cleaned.strip():
                                all_text.append(cleaned)

        except Exception as e:
            logger.error(f"[PPTX] Error processing slides in {file_path}: {e}")

        return self.create_sample(all_text, embedded_images, DocumentMetadata(file_path=file_path))
